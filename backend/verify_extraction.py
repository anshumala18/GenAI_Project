
import io
import pdfplumber
import docx
from pptx import Presentation
from processor import DocumentProcessor

def test_extractors():
    processor = DocumentProcessor()
    
    # 1. Test DOCX
    doc = docx.Document()
    doc.add_paragraph("Hello World from DOCX")
    doc_io = io.BytesIO()
    doc.save(doc_io)
    doc_content = doc_io.getvalue()
    
    doc_text = processor.extract_text_from_docx(doc_content)
    print(f"DOCX Extraction: {'SUCCESS' if 'Hello World' in doc_text else 'FAILED'}")
    
    # 2. Test PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Hello World from PPTX"
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_content = ppt_io.getvalue()
    
    ppt_text = processor.extract_text_from_pptx(ppt_content)
    print(f"PPTX Extraction: {'SUCCESS' if 'Hello World' in ppt_text else 'FAILED'}")

if __name__ == "__main__":
    test_extractors()
