import pdfplumber
import docx
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
import io

class DocumentProcessor:
    def __init__(self, chunk_size=1500, chunk_overlap=200):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
        )

    def extract_text_from_pdf(self, file_content: bytes) -> str:
        text = ""
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text

    def extract_text_from_docx(self, file_content: bytes) -> str:
        doc = docx.Document(io.BytesIO(file_content))
        text = "\n".join([para.text for para in doc.paragraphs])
        return text

    def extract_text_from_pptx(self, file_content: bytes) -> str:
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def clean_text(self, text: str) -> str:
        # Basic cleaning: reduce multiple newlines and spaces
        import re
        text = re.sub(r'\n+', '\n', text)
        text = re.sub(r' +', ' ', text)
        return text.strip()

    def get_chunks(self, text: str):
        return self.text_splitter.split_text(text)
